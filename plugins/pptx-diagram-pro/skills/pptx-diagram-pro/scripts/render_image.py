"""Render a diagram IR (JSON) to a PPTX by generating a single diagram image via Gemini.

Generates a 4K-style diagram image that visualizes the IR, then inserts it as a
full-bleed slide in a new PPTX file. Use the native renderer for editable shapes.

Requires:
    pip install python-pptx google-genai
    GEMINI_API_KEY environment variable

Usage:
    python render_image.py <ir.json> [--output <out.pptx>] [--png-only]

Examples:
    python render_image.py arch.json
    python render_image.py flow.json --output ./out/flow.pptx
    python render_image.py flow.json --png-only      # only save PNG, no PPTX
"""

import argparse
import base64
import json
import os
import sys
import time

try:
    from pptx import Presentation
    from pptx.util import Inches, Emu
except ImportError:
    print("ERROR: python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


PALETTE = {
    "technical-report": {"primary": "#1E3A5F", "secondary": "#4A90A4", "accent": "#E07B39", "bg": "#F5F7FA"},
    "clarity":          {"primary": "#2C2C2C", "secondary": "#6B6B6B", "accent": "#D4A017", "bg": "#FAFAFA"},
    "tech-focus":       {"primary": "#0B3D91", "secondary": "#3F8EFC", "accent": "#FF6B35", "bg": "#EEF3FB"},
    "growth":           {"primary": "#1B5E20", "secondary": "#66BB6A", "accent": "#FFA726", "bg": "#F1F8E9"},
}

KIND_VERBAL = {
    "start": "라운드 사각형 (시작)",
    "end": "라운드 사각형 (종료)",
    "process": "직사각형 (처리 단계)",
    "decision": "마름모 (분기/판단)",
    "io": "평행사변형 (입출력)",
    "service": "라운드 사각형 (서비스/모듈)",
    "database": "원통형 (데이터베이스)",
    "queue": "양단 직선 평행사변형 (메시지 큐)",
    "model": "육각형 (AI/ML 모델)",
    "user": "사람 형상 또는 라운드 사각형 (사용자)",
    "external": "점선 외곽선 라운드 사각형 (외부 시스템)",
    "note": "직각 사각형 (메모)",
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def build_prompt(ir: dict) -> str:
    mood = ir.get("style", {}).get("mood", "technical-report")
    palette = PALETTE.get(mood, PALETTE["technical-report"])
    layout = ir["layout"]
    direction_label = "좌→우" if layout.get("direction") == "LR" else "위→아래"

    lines = []
    lines.append("■ 이미지 유형")
    lines.append(
        "PPT 슬라이드 스타일의 다이어그램. 정부 국가연구개발 과제 발표자료 수준의 전문적 도식. "
        "가로형(16:9), 4K 고해상도, 플랫 인포그래픽, 벡터 그래픽 스타일."
    )
    lines.append("")
    lines.append("■ 색상 팔레트 (4색 제한)")
    lines.append(f"1. 주조색: {palette['primary']} — 제목, 노드 외곽선, 강조 텍스트")
    lines.append(f"2. 보조색: {palette['secondary']} — 연결선, 그룹 외곽")
    lines.append(f"3. 강조색: {palette['accent']} — 핵심 노드 채움")
    lines.append(f"4. 배경색: {palette['bg']} — 캔버스, 그룹 채움")
    lines.append("")
    lines.append("■ 다이어그램 유형 및 흐름")
    lines.append(f"- 유형: {ir['type']} (flowchart=프로세스 흐름, architecture=시스템 구성)")
    lines.append(f"- 격자: {layout['rows']}행 × {layout['cols']}열")
    lines.append(f"- 흐름 방향: {direction_label}")
    lines.append("")
    lines.append("■ 노드 구성 (격자 좌표 row/col은 0부터 시작, 좌상단=0,0)")
    for n in ir["nodes"]:
        accent = " [강조: 채움색=강조색, 흰색 라벨]" if n.get("accent") else ""
        rs = f", rowspan={n['rowspan']}" if n.get("rowspan", 1) > 1 else ""
        cs = f", colspan={n['colspan']}" if n.get("colspan", 1) > 1 else ""
        kind = KIND_VERBAL.get(n["kind"], n["kind"])
        lines.append(f"- ({n['row']},{n['col']}{rs}{cs}) {kind}: \"{n['label']}\"{accent}")
    lines.append("")

    if ir.get("groups"):
        lines.append("■ 그룹 (회색 배경 라운드 사각형으로 묶음)")
        for g in ir["groups"]:
            lines.append(
                f"- ({g['row']},{g['col']}, rowspan={g.get('rowspan', 1)}, colspan={g.get('colspan', 1)}) \"{g.get('label', '')}\""
            )
        lines.append("")

    lines.append("■ 연결선 (화살표)")
    for e in ir.get("edges", []) or []:
        style = e.get("style", "solid")
        style_label = "실선" if style == "solid" else ("점선" if style == "dashed" else "도트선")
        lab = f" 라벨=\"{e['label']}\"" if e.get("label") else ""
        lines.append(f"- {e['from']} → {e['to']} ({style_label}){lab}")
    lines.append("")

    lines.append("■ 반드시 포함할 텍스트 (한글, 정확히 렌더링)")
    lines.append(f"- 슬라이드 제목: \"{ir.get('title', '')}\"")
    if ir.get("subtitle"):
        lines.append(f"- 부제: \"{ir['subtitle']}\"")
    for n in ir["nodes"]:
        lines.append(f"- 노드 라벨: \"{n['label']}\"")
    for e in ir.get("edges", []) or []:
        if e.get("label"):
            lines.append(f"- 엣지 라벨: \"{e['label']}\"")
    lines.append("")

    lines.append("■ 스타일 지침")
    lines.append("- 모든 텍스트 한글 렌더링, 깔끔하고 선명하게")
    lines.append("- 정부 공식 문서 수준의 격식, 직각 또는 살짝 라운드 도형, 명확한 테두리")
    lines.append("- 화살표는 보조색으로 1~1.5pt 두께, 머리 폭 적정")
    lines.append("- 한·영 병기 금지 (고유 약어 AI, ML, VLM, REST API 등은 영문만 허용)")
    lines.append("- 그라데이션/그림자/글래스모피즘 금지")
    lines.append("- 이모지·이미지 플레이스홀더(`[Icon]`, `[사진]`) 금지")
    lines.append("- 렌더링 결과에 워터마크·placeholder bracket 절대 포함 금지")
    return "\n".join(lines)


def call_gemini(prompt: str, output_png: str, model: str = "gemini-2.5-flash-image-preview", retries: int = 3):
    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        print("ERROR: GEMINI_API_KEY environment variable is not set.", file=sys.stderr)
        sys.exit(1)

    try:
        from google import genai
    except ImportError:
        print("ERROR: google-genai is required for image mode. Install with: pip install google-genai", file=sys.stderr)
        sys.exit(1)

    client = genai.Client(api_key=api_key)

    last_err = None
    for attempt in range(1, retries + 1):
        try:
            response = client.models.generate_content(model=model, contents=[prompt])
            for cand in response.candidates or []:
                for part in (cand.content.parts or []):
                    inline = getattr(part, "inline_data", None)
                    if inline and inline.data:
                        data = inline.data
                        if isinstance(data, str):
                            data = base64.b64decode(data)
                        with open(output_png, "wb") as f:
                            f.write(data)
                        return output_png
            last_err = "no inline image data in response"
        except Exception as e:
            last_err = str(e)
        if attempt < retries:
            time.sleep(5)
    print(f"ERROR: Gemini image generation failed after {retries} attempts: {last_err}", file=sys.stderr)
    sys.exit(3)


def build_pptx(image_path: str, ir: dict, output_path: str) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(image_path, left=Emu(0), top=Emu(0), width=SLIDE_W, height=SLIDE_H)

    notes = ir.get("title", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    prs.core_properties.title = ir.get("title", "Diagram")
    prs.save(output_path)
    return output_path


def render(ir_path: str, output_path: str | None = None, png_only: bool = False) -> str:
    with open(ir_path, "r", encoding="utf-8") as f:
        ir = json.load(f)

    base = os.path.splitext(os.path.basename(ir_path))[0]
    out_dir = os.path.dirname(os.path.abspath(ir_path)) or "."
    png_path = os.path.join(out_dir, f"{base}.png")

    prompt = build_prompt(ir)
    call_gemini(prompt, png_path)

    if png_only:
        return png_path

    if output_path is None:
        output_path = os.path.join(out_dir, f"{base}.pptx")
    build_pptx(png_path, ir, output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Render diagram IR JSON to PPTX via Gemini image generation")
    parser.add_argument("ir_path", help="Path to IR JSON file")
    parser.add_argument("--output", dest="output_path", help="Output PPTX path")
    parser.add_argument("--png-only", action="store_true", help="Only emit PNG, skip PPTX wrapping")
    args = parser.parse_args()
    if not os.path.isfile(args.ir_path):
        print(f"ERROR: IR file not found: {args.ir_path}", file=sys.stderr)
        sys.exit(1)
    out = render(args.ir_path, args.output_path, args.png_only)
    kind = "PNG" if args.png_only else "PPTX"
    print(f"{kind} created: {out}")


if __name__ == "__main__":
    main()
