"""Render a diagram IR (JSON) to a single-slide PPTX using python-pptx native shapes.

The IR schema is documented in ../references/ir-spec.md.

Usage:
    python render_native.py <ir.json> [--output <out.pptx>] [--title <override-title>]

Examples:
    python render_native.py flow.json
    python render_native.py arch.json --output ./out/arch.pptx
"""

import argparse
import json
import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
except ImportError:
    print("ERROR: python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

try:
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    HAS_DASH = True
except ImportError:
    HAS_DASH = False

from lxml import etree


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TITLE_H = Inches(1.0)
MARGIN = Inches(0.6)
BODY_TOP = TITLE_H + Inches(0.2)

PALETTE = {
    "technical-report": {"primary": "#1E3A5F", "secondary": "#4A90A4", "accent": "#E07B39", "bg": "#F5F7FA"},
    "clarity":          {"primary": "#2C2C2C", "secondary": "#6B6B6B", "accent": "#D4A017", "bg": "#FAFAFA"},
    "tech-focus":       {"primary": "#0B3D91", "secondary": "#3F8EFC", "accent": "#FF6B35", "bg": "#EEF3FB"},
    "growth":           {"primary": "#1B5E20", "secondary": "#66BB6A", "accent": "#FFA726", "bg": "#F1F8E9"},
}

KIND_SHAPE = {
    "start": MSO_SHAPE.ROUNDED_RECTANGLE,
    "end": MSO_SHAPE.ROUNDED_RECTANGLE,
    "process": MSO_SHAPE.RECTANGLE,
    "decision": MSO_SHAPE.DIAMOND,
    "io": MSO_SHAPE.PARALLELOGRAM,
    "service": MSO_SHAPE.ROUNDED_RECTANGLE,
    "database": MSO_SHAPE.CAN,
    "queue": MSO_SHAPE.PARALLELOGRAM,
    "model": MSO_SHAPE.HEXAGON,
    "user": MSO_SHAPE.ROUNDED_RECTANGLE,
    "external": MSO_SHAPE.ROUNDED_RECTANGLE,
    "note": MSO_SHAPE.RECTANGLE,
}


def hex_to_rgb(s: str) -> RGBColor:
    s = s.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def lighten(hex_color: str, factor: float = 0.85) -> RGBColor:
    s = hex_color.lstrip("#")
    r, g, b = int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return RGBColor(r, g, b)


def set_arrow(connector, head_type: str = "triangle"):
    """Add a tail arrow head to a connector via direct XML."""
    ln = connector.line._get_or_add_ln()
    existing = ln.find(qn("a:tailEnd"))
    if existing is not None:
        ln.remove(existing)
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", head_type)
    tail.set("w", "med")
    tail.set("len", "med")


def set_dash(connector, style: str):
    """Apply dash style to a connector line."""
    if not HAS_DASH:
        return
    if style == "dashed":
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    elif style == "dotted":
        connector.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT


def grid_cell(layout, body_left, body_top, body_w, body_h, row, col, rowspan=1, colspan=1, padding=Inches(0.1)):
    cell_w = body_w / layout["cols"]
    cell_h = body_h / layout["rows"]
    left = body_left + cell_w * col + padding
    top = body_top + cell_h * row + padding
    width = cell_w * colspan - 2 * padding
    height = cell_h * rowspan - 2 * padding
    return int(left), int(top), int(width), int(height)


def set_text(shape, text, font_size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), font_name="Malgun Gothic"):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.text = ""
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    for i, line in enumerate(str(text).split("\n")):
        if i == 0:
            run = para.add_run()
        else:
            new_para = tf.add_paragraph()
            new_para.alignment = PP_ALIGN.CENTER
            run = new_para.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


def set_dashed_outline(shape):
    ln = shape.line._get_or_add_ln()
    existing = ln.find(qn("a:prstDash"))
    if existing is not None:
        ln.remove(existing)
    prst = etree.SubElement(ln, qn("a:prstDash"))
    prst.set("val", "dash")


def add_node(slide, node, layout, palette, body_box):
    body_left, body_top, body_w, body_h = body_box
    left, top, width, height = grid_cell(
        layout, body_left, body_top, body_w, body_h,
        node["row"], node["col"], node.get("rowspan", 1), node.get("colspan", 1),
    )
    shape_type = KIND_SHAPE.get(node["kind"], MSO_SHAPE.RECTANGLE)
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)

    accent = bool(node.get("accent"))
    if accent:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(palette["accent"])
        text_color = RGBColor(0xFF, 0xFF, 0xFF)
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = lighten(palette["secondary"], 0.85)
        text_color = hex_to_rgb(palette["primary"])

    shape.line.color.rgb = hex_to_rgb(palette["primary"])
    shape.line.width = Pt(1.5)

    if node["kind"] == "external":
        set_dashed_outline(shape)

    set_text(shape, node["label"], font_size=14, bold=True, color=text_color)
    return shape


def add_group(slide, group, layout, palette, body_box):
    body_left, body_top, body_w, body_h = body_box
    left, top, width, height = grid_cell(
        layout, body_left, body_top, body_w, body_h,
        group["row"], group["col"], group.get("rowspan", 1), group.get("colspan", 1),
        padding=Inches(0.04),
    )
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(palette["bg"])
    shape.line.color.rgb = hex_to_rgb(palette["secondary"])
    shape.line.width = Pt(0.75)

    label = group.get("label")
    if label:
        from pptx.util import Emu as _Emu
        tb = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.05), width - Inches(0.2), Inches(0.3))
        tf = tb.text_frame
        tf.margin_left = Inches(0.02)
        tf.margin_top = Inches(0.0)
        tf.text = ""
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = label
        run.font.name = "Malgun Gothic"
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = hex_to_rgb(palette["secondary"])
    return shape


def anchor_point(shape, side: str):
    if side == "right":
        return shape.left + shape.width, shape.top + shape.height // 2
    if side == "left":
        return shape.left, shape.top + shape.height // 2
    if side == "bottom":
        return shape.left + shape.width // 2, shape.top + shape.height
    if side == "top":
        return shape.left + shape.width // 2, shape.top
    return shape.left + shape.width // 2, shape.top + shape.height // 2


def choose_sides(from_node, to_node):
    dr = to_node["row"] - from_node["row"]
    dc = to_node["col"] - from_node["col"]
    if dr == 0 and dc != 0:
        return ("right", "left") if dc > 0 else ("left", "right")
    if dc == 0 and dr != 0:
        return ("bottom", "top") if dr > 0 else ("top", "bottom")
    if dc != 0 and dr != 0:
        # Diagonal: prefer horizontal-then-vertical elbow
        return ("right", "top") if (dc > 0 and dr > 0) else \
               ("right", "bottom") if (dc > 0 and dr < 0) else \
               ("left", "top") if (dc < 0 and dr > 0) else \
               ("left", "bottom")
    return ("right", "left")


def add_edge(slide, edge, node_shapes, node_specs, palette):
    from_id, to_id = edge["from"], edge["to"]
    if from_id not in node_shapes or to_id not in node_shapes:
        return None
    from_shape = node_shapes[from_id]
    to_shape = node_shapes[to_id]
    from_spec = node_specs[from_id]
    to_spec = node_specs[to_id]

    from_side, to_side = choose_sides(from_spec, to_spec)
    bend = edge.get("bend", "auto")
    use_elbow = (from_spec["row"] != to_spec["row"]) and (from_spec["col"] != to_spec["col"])
    if bend == "horizontal" or bend == "vertical":
        use_elbow = True
    connector_type = MSO_CONNECTOR.ELBOW if use_elbow else MSO_CONNECTOR.STRAIGHT

    bx, by = anchor_point(from_shape, from_side)
    ex, ey = anchor_point(to_shape, to_side)

    connector = slide.shapes.add_connector(connector_type, bx, by, ex, ey)
    connector.line.color.rgb = hex_to_rgb(palette["secondary"])
    connector.line.width = Pt(1.25)

    set_dash(connector, edge.get("style", "solid"))
    if edge.get("arrow", "single") != "none":
        set_arrow(connector, "triangle")

    label = edge.get("label")
    if label:
        mx, my = (bx + ex) // 2, (by + ey) // 2
        tb_w = Inches(1.0)
        tb_h = Inches(0.3)
        tb = slide.shapes.add_textbox(mx - tb_w // 2, my - tb_h // 2, tb_w, tb_h)
        tb.fill.solid()
        tb.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tb.line.fill.background()
        tf = tb.text_frame
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.0)
        tf.margin_bottom = Inches(0.0)
        tf.text = ""
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.name = "Malgun Gothic"
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0x3A, 0x3A, 0x3A)

    return connector


def add_title(slide, title, subtitle, palette):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(palette["primary"])
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(MARGIN, Inches(0.18), SLIDE_W - 2 * MARGIN, Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = hex_to_rgb(palette["primary"])

    if subtitle:
        sub_p = tf.add_paragraph()
        sub_p.alignment = PP_ALIGN.LEFT
        sr = sub_p.add_run()
        sr.text = subtitle
        sr.font.name = "Malgun Gothic"
        sr.font.size = Pt(16)
        sr.font.color.rgb = hex_to_rgb(palette["secondary"])


def validate_ir(ir):
    errs = []
    if "type" not in ir or ir["type"] not in ("flowchart", "architecture"):
        errs.append("ir.type must be 'flowchart' or 'architecture'")
    if "nodes" not in ir or not ir["nodes"]:
        errs.append("ir.nodes must be a non-empty list")
    if "layout" not in ir or "rows" not in ir["layout"] or "cols" not in ir["layout"]:
        errs.append("ir.layout.rows and .cols are required")
    ids = {n["id"] for n in ir.get("nodes", [])}
    for e in ir.get("edges", []):
        if e.get("from") not in ids or e.get("to") not in ids:
            errs.append(f"edge references unknown node id: {e}")
    accent_count = sum(1 for n in ir.get("nodes", []) if n.get("accent"))
    if accent_count > 3:
        errs.append(f"too many accent nodes: {accent_count} (max 3)")
    return errs


def render(ir_path: str, output_path: str | None = None, title_override: str | None = None) -> str:
    with open(ir_path, "r", encoding="utf-8") as f:
        ir = json.load(f)

    errs = validate_ir(ir)
    if errs:
        for e in errs:
            print(f"IR ERROR: {e}", file=sys.stderr)
        sys.exit(2)

    style = ir.get("style", {})
    mood = style.get("mood", "technical-report")
    palette = PALETTE.get(mood, PALETTE["technical-report"])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, title_override or ir.get("title", ""), ir.get("subtitle"), palette)

    body_left = MARGIN
    body_top = BODY_TOP
    body_w = SLIDE_W - 2 * MARGIN
    body_h = SLIDE_H - BODY_TOP - Inches(0.4)
    body_box = (body_left, body_top, body_w, body_h)

    for group in ir.get("groups", []) or []:
        add_group(slide, group, ir["layout"], palette, body_box)

    node_shapes = {}
    node_specs = {}
    for node in ir["nodes"]:
        node_specs[node["id"]] = node
        node_shapes[node["id"]] = add_node(slide, node, ir["layout"], palette, body_box)

    for edge in ir.get("edges", []) or []:
        add_edge(slide, edge, node_shapes, node_specs, palette)

    if output_path is None:
        base = os.path.splitext(os.path.basename(ir_path))[0]
        output_path = os.path.join(os.path.dirname(os.path.abspath(ir_path)) or ".", f"{base}.pptx")

    prs.core_properties.title = ir.get("title", "Diagram")
    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Render diagram IR JSON to PPTX (native python-pptx shapes)")
    parser.add_argument("ir_path", help="Path to IR JSON file")
    parser.add_argument("--output", dest="output_path", help="Output PPTX path")
    parser.add_argument("--title", dest="title_override", help="Override the diagram title")
    args = parser.parse_args()
    if not os.path.isfile(args.ir_path):
        print(f"ERROR: IR file not found: {args.ir_path}", file=sys.stderr)
        sys.exit(1)
    out = render(args.ir_path, args.output_path, args.title_override)
    print(f"PPTX created: {out}")


if __name__ == "__main__":
    main()
