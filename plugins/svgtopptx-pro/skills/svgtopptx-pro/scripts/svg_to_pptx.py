"""
SVG to editable PPTX converter.
Parses SVG elements (rect, text, circle, line, path, polygon) and recreates them
as native PowerPoint shapes so they remain fully editable.

Usage:
    python svg_to_pptx.py <input_dir_or_file> [output.pptx]

    - input_dir_or_file: Directory containing SVG files, or a single SVG file path
    - output.pptx: Optional output path (default: <input_dir>/output.pptx)

Examples:
    python svg_to_pptx.py ./images/
    python svg_to_pptx.py ./images/ ./result.pptx
    python svg_to_pptx.py ./diagram.svg
    python svg_to_pptx.py ./diagram.svg ./diagram.pptx
"""
import re
import os
import sys
import glob
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn


# ─────────────────────────── Color / Style Parsing ───────────────────────────

def parse_color(color_str):
    """Parse color from hex (#RRGGBB) or rgb(r,g,b) or rgba() format."""
    if not color_str or color_str == 'none':
        return None
    color_str = color_str.strip()
    if color_str.startswith('#'):
        h = color_str.lstrip('#')
        if len(h) == 3:
            h = ''.join(c * 2 for c in h)
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    m = re.match(r'rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)', color_str)
    if m:
        return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    return None


def parse_style(style_str):
    """Parse inline CSS style string into dict."""
    result = {}
    if not style_str:
        return result
    for part in style_str.split(';'):
        part = part.strip()
        if ':' in part:
            k, v = part.split(':', 1)
            result[k.strip()] = v.strip()
    return result


def get_fill(elem, style):
    fill_str = style.get('fill') or elem.get('fill')
    return parse_color(fill_str)


def get_stroke(elem, style):
    stroke_str = style.get('stroke') or elem.get('stroke')
    return parse_color(stroke_str)


def get_stroke_width(elem, style):
    sw = style.get('stroke-width') or elem.get('stroke-width', '1')
    sw = sw.replace('px', '').strip()
    try:
        return float(sw)
    except ValueError:
        return 1.0


def get_font_size(style):
    fs = style.get('font-size', '12px').replace('px', '').strip()
    try:
        return float(fs)
    except ValueError:
        return 12.0


def get_font_weight(style):
    fw = style.get('font-weight', '400')
    try:
        return int(fw) >= 500
    except ValueError:
        return fw == 'bold'


# ─────────────────────────── Element Helpers ─────────────────────────────────

def get_local_tag(elem):
    """Get local tag name handling namespaces and comment nodes."""
    tag = elem.tag
    if callable(tag):
        return ''
    tag = str(tag)
    if '}' in tag:
        return etree.QName(tag).localname
    return tag


def get_text_content(elem):
    """Get all text content from element and tspan children."""
    text = elem.text or ''
    for child in elem:
        t = get_local_tag(child)
        if t == 'tspan':
            text += child.text or ''
        text += child.tail or ''
    return text.strip()


def collect_elements(parent, elements):
    """Recursively collect drawable SVG elements."""
    for child in parent:
        tag = get_local_tag(child)
        if not tag or tag == 'defs':
            continue
        if tag == 'g':
            collect_elements(child, elements)
        elif tag in ('rect', 'text', 'circle', 'line', 'polygon', 'path'):
            elements.append(child)


def set_arrow_end(connector):
    """Add triangle arrow head to connector end."""
    try:
        cxnSp = connector._element
        spPr = cxnSp.find(qn('p:spPr'))
        ln = spPr.find(qn('a:ln')) if spPr is not None else None
        if ln is not None:
            tail = ln.find(qn('a:tailEnd'))
            if tail is None:
                tail = etree.SubElement(ln, qn('a:tailEnd'))
            tail.set('type', 'triangle')
            tail.set('w', 'med')
            tail.set('len', 'med')
    except Exception:
        pass


# ─────────────────────────── Shape Builders ──────────────────────────────────

def add_rect(slide, elem, style, s, off_x, off_y):
    x, y = float(elem.get('x', 0)), float(elem.get('y', 0))
    w, h = float(elem.get('width', 0)), float(elem.get('height', 0))

    left, top = int(x * s + off_x), int(y * s + off_y)
    width, height = int(w * s), int(h * s)

    if width < Emu(10000) or height < Emu(10000):
        return

    shape = slide.shapes.add_shape(5, left, top, width, height)  # ROUNDED_RECTANGLE

    fill_color = get_fill(elem, style)
    stroke_color = get_stroke(elem, style)
    sw = get_stroke_width(elem, style)

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if stroke_color and style.get('stroke') != 'none':
        shape.line.color.rgb = stroke_color
        shape.line.width = Pt(sw)
    else:
        shape.line.fill.background()


def add_text(slide, elem, style, s, off_x, off_y):
    x, y = float(elem.get('x', 0)), float(elem.get('y', 0))
    text_content = get_text_content(elem)
    if not text_content:
        return

    font_size = get_font_size(style)
    is_bold = get_font_weight(style)
    fill_color = get_fill(elem, style)
    text_anchor = style.get('text-anchor', elem.get('text-anchor', 'start'))

    char_w = font_size * 0.7
    est_w = max(len(text_content) * char_w, 60)
    est_h = font_size * 1.8

    if text_anchor == 'middle':
        left = int((x - est_w / 2) * s + off_x)
    elif text_anchor == 'end':
        left = int((x - est_w) * s + off_x)
    else:
        left = int(x * s + off_x)

    top = int((y - font_size) * s + off_y)
    width = max(int(est_w * s), Emu(100000))
    height = max(int(est_h * s), Emu(50000))

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False

    p = tf.paragraphs[0]
    align_map = {'middle': PP_ALIGN.CENTER, 'end': PP_ALIGN.RIGHT}
    p.alignment = align_map.get(text_anchor, PP_ALIGN.LEFT)

    run = p.add_run()
    run.text = text_content
    run.font.size = Pt(font_size * 0.75)
    run.font.bold = is_bold
    run.font.name = 'Malgun Gothic'
    if fill_color:
        run.font.color.rgb = fill_color


def add_circle(slide, elem, style, s, off_x, off_y):
    cx, cy = float(elem.get('cx', 0)), float(elem.get('cy', 0))
    r = float(elem.get('r', 0))
    diameter = int(2 * r * s)
    if diameter < Emu(10000):
        return

    left, top = int((cx - r) * s + off_x), int((cy - r) * s + off_y)
    shape = slide.shapes.add_shape(9, left, top, diameter, diameter)  # OVAL

    fill_color = get_fill(elem, style)
    stroke_color = get_stroke(elem, style)
    sw = get_stroke_width(elem, style)

    if fill_color and style.get('fill') != 'none' and elem.get('fill') != 'none':
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if stroke_color and style.get('stroke') != 'none':
        shape.line.color.rgb = stroke_color
        shape.line.width = Pt(sw)
        dash = style.get('stroke-dasharray') or elem.get('stroke-dasharray')
        if dash and dash != 'none':
            shape.line.dash_style = 3  # DASH
    else:
        shape.line.fill.background()


def add_line(slide, elem, style, s, off_x, off_y):
    x1, y1 = float(elem.get('x1', 0)), float(elem.get('y1', 0))
    x2, y2 = float(elem.get('x2', 0)), float(elem.get('y2', 0))

    conn = slide.shapes.add_connector(
        1, int(x1 * s + off_x), int(y1 * s + off_y),
        int(x2 * s + off_x), int(y2 * s + off_y))

    stroke_color = get_stroke(elem, style)
    sw = get_stroke_width(elem, style)
    if stroke_color:
        conn.line.color.rgb = stroke_color
    conn.line.width = Pt(sw)

    dash = style.get('stroke-dasharray') or elem.get('stroke-dasharray')
    if dash and dash != 'none':
        conn.line.dash_style = 3

    if elem.get('marker-end') and 'arrow' in elem.get('marker-end', ''):
        set_arrow_end(conn)


def add_polygon(slide, elem, style, s, off_x, off_y):
    points_str = elem.get('points', '').strip()
    if not points_str:
        return

    coords = []
    for pair in points_str.split():
        parts = pair.split(',')
        if len(parts) == 2:
            coords.append((float(parts[0]), float(parts[1])))
    if len(coords) < 3:
        return

    xs = [c[0] for c in coords]
    ys = [c[1] for c in coords]
    min_x, max_x = min(xs), max(xs)
    min_y, max_y = min(ys), max(ys)
    w, h = max(max_x - min_x, 1), max(max_y - min_y, 1)

    left, top = int(min_x * s + off_x), int(min_y * s + off_y)
    width = max(int(w * s), Emu(5000))
    height = max(int(h * s), Emu(5000))

    shape = slide.shapes.add_shape(7, left, top, width, height)  # ISOSCELES_TRIANGLE
    fill_color = get_fill(elem, style)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    stroke_color = get_stroke(elem, style)
    if stroke_color and style.get('stroke') != 'none':
        shape.line.color.rgb = stroke_color
    else:
        shape.line.fill.background()


def parse_path_points(d):
    """Extract key points from SVG path data (M/L/Q/C/Z commands)."""
    points = []
    tokens = re.findall(r'[MLQCZmlqcz]|[-+]?\d*\.?\d+', d)
    i, cx, cy = 0, 0, 0

    while i < len(tokens):
        cmd = tokens[i]
        if cmd in ('M', 'm'):
            i += 1
            if i + 1 < len(tokens):
                x, y = float(tokens[i]), float(tokens[i + 1])
                if cmd == 'm':
                    x += cx; y += cy
                cx, cy = x, y
                points.append((x, y))
                i += 2
        elif cmd in ('L', 'l'):
            i += 1
            if i + 1 < len(tokens):
                x, y = float(tokens[i]), float(tokens[i + 1])
                if cmd == 'l':
                    x += cx; y += cy
                cx, cy = x, y
                points.append((x, y))
                i += 2
        elif cmd in ('Q', 'q'):
            i += 1
            if i + 3 < len(tokens):
                i += 2  # skip control point
                x, y = float(tokens[i]), float(tokens[i + 1])
                if cmd == 'q':
                    x += cx; y += cy
                cx, cy = x, y
                points.append((x, y))
                i += 2
        elif cmd in ('C', 'c'):
            i += 1
            if i + 5 < len(tokens):
                i += 4  # skip control points
                x, y = float(tokens[i]), float(tokens[i + 1])
                if cmd == 'c':
                    x += cx; y += cy
                cx, cy = x, y
                points.append((x, y))
                i += 2
        elif cmd in ('Z', 'z'):
            i += 1
        else:
            try:
                x = float(tokens[i])
                if i + 1 < len(tokens):
                    y = float(tokens[i + 1])
                    cx, cy = x, y
                    points.append((x, y))
                    i += 2
                else:
                    i += 1
            except ValueError:
                i += 1
    return points


def add_path(slide, elem, style, s, off_x, off_y):
    """Approximate SVG <path> as connected line segments."""
    d = elem.get('d', '').strip()
    if not d:
        return

    stroke_color = get_stroke(elem, style)
    sw = get_stroke_width(elem, style)
    points = parse_path_points(d)
    if len(points) < 2:
        return

    for i in range(len(points) - 1):
        x1, y1 = points[i]
        x2, y2 = points[i + 1]
        try:
            conn = slide.shapes.add_connector(
                1, int(x1 * s + off_x), int(y1 * s + off_y),
                int(x2 * s + off_x), int(y2 * s + off_y))
            if stroke_color:
                conn.line.color.rgb = stroke_color
            conn.line.width = Pt(sw)

            dash = style.get('stroke-dasharray') or elem.get('stroke-dasharray')
            if dash and dash != 'none':
                conn.line.dash_style = 3

            if i == len(points) - 2:
                if elem.get('marker-end') and 'arrow' in elem.get('marker-end', ''):
                    set_arrow_end(conn)
        except Exception:
            pass


# ─────────────────────────── Main Conversion ─────────────────────────────────

SHAPE_HANDLERS = {
    'rect': add_rect,
    'text': add_text,
    'circle': add_circle,
    'line': add_line,
    'polygon': add_polygon,
    'path': add_path,
}


def process_svg(svg_path, prs):
    """Parse one SVG file and add it as a new slide in the presentation."""
    tree = etree.parse(svg_path)
    root = tree.getroot()

    viewbox = root.get('viewBox', '0 0 680 480')
    parts = viewbox.split()
    vb_w, vb_h = float(parts[2]), float(parts[3])

    s = min(prs.slide_width / vb_w, prs.slide_height / vb_h)
    off_x = (prs.slide_width - vb_w * s) / 2
    off_y = (prs.slide_height - vb_h * s) / 2

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    elements = []
    collect_elements(root, elements)

    for elem in elements:
        tag = get_local_tag(elem)
        style = parse_style(elem.get('style', ''))
        handler = SHAPE_HANDLERS.get(tag)
        if handler:
            handler(slide, elem, style, s, off_x, off_y)

    return slide


def convert(svg_files, output_path):
    """Convert a list of SVG files into a single editable PPTX."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for svg_file in svg_files:
        fname = os.path.basename(svg_file)
        print(f"  [{fname}]", end=' ')
        slide = process_svg(svg_file, prs)
        print(f"-> {len(slide.shapes)} shapes")

    prs.save(output_path)
    print(f"\nSaved: {output_path}  ({len(prs.slides)} slides)")
    return output_path


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    input_path = os.path.abspath(sys.argv[1])
    output_path = sys.argv[2] if len(sys.argv) >= 3 else None

    if os.path.isdir(input_path):
        svg_files = sorted(glob.glob(os.path.join(input_path, '*.svg')))
        if not output_path:
            output_path = os.path.join(input_path, 'output.pptx')
    elif os.path.isfile(input_path) and input_path.lower().endswith('.svg'):
        svg_files = [input_path]
        if not output_path:
            output_path = os.path.splitext(input_path)[0] + '.pptx'
    else:
        print(f"Error: '{input_path}' is not a valid SVG file or directory.")
        sys.exit(1)

    if not svg_files:
        print(f"No SVG files found in '{input_path}'.")
        sys.exit(1)

    print(f"Found {len(svg_files)} SVG file(s):")
    convert(svg_files, output_path)


if __name__ == '__main__':
    main()
