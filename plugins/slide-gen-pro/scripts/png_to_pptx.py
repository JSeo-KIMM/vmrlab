"""Convert slide PNG images to a PPTX presentation.

Each PNG is inserted as a full-bleed slide image (16:9).
Optionally reads prompt markdown files to extract slide titles for speaker notes.

Usage:
    python png_to_pptx.py <images_dir> [--prompts <prompts_dir>] [--output <output.pptx>] [--title <title>]

Examples:
    python png_to_pptx.py ./slides/images
    python png_to_pptx.py ./slides/images --prompts ./slides/prompts --output presentation.pptx
    python png_to_pptx.py ./slides/images --title "연구개발 발표자료" --output result.pptx
"""

import argparse
import glob
import os
import re
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Emu
except ImportError:
    print("ERROR: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def extract_slide_title(prompt_path: str) -> str:
    """Extract the first heading or '반드시 포함할 텍스트' section from a prompt markdown."""
    try:
        with open(prompt_path, "r", encoding="utf-8") as f:
            content = f.read()
    except Exception:
        return ""

    heading = re.search(r"^#\s+(.+)$", content, re.MULTILINE)
    if heading:
        return heading.group(1).strip()

    title_section = re.search(r"반드시 포함할 텍스트.*?\n-\s*[\""](.+?)[\""]", content, re.DOTALL)
    if title_section:
        return title_section.group(1).strip()

    return ""


def find_matching_prompt(image_name: str, prompts_dir: str) -> str | None:
    """Find a prompt file that matches the image filename stem."""
    stem = os.path.splitext(image_name)[0]
    for ext in (".md", ".txt"):
        candidate = os.path.join(prompts_dir, stem + ext)
        if os.path.isfile(candidate):
            return candidate
    return None


def build_pptx(
    images_dir: str,
    prompts_dir: str | None = None,
    output_path: str | None = None,
    title: str = "Presentation",
) -> str:
    """Build a PPTX from PNG images.

    Returns the path to the generated PPTX file.
    """
    image_files = sorted(
        glob.glob(os.path.join(images_dir, "*.png"))
        + glob.glob(os.path.join(images_dir, "*.jpg"))
        + glob.glob(os.path.join(images_dir, "*.jpeg"))
    )

    if not image_files:
        print(f"ERROR: No image files found in {images_dir}")
        sys.exit(1)

    if output_path is None:
        parent = os.path.dirname(os.path.abspath(images_dir))
        output_path = os.path.join(parent, "presentation.pptx")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]

    for img_path in image_files:
        slide = prs.slides.add_slide(blank_layout)

        slide.shapes.add_picture(
            img_path,
            left=Emu(0),
            top=Emu(0),
            width=SLIDE_WIDTH,
            height=SLIDE_HEIGHT,
        )

        if prompts_dir:
            prompt_file = find_matching_prompt(os.path.basename(img_path), prompts_dir)
            if prompt_file:
                slide_title = extract_slide_title(prompt_file)
                if slide_title:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = slide_title

    if title:
        prs.core_properties.title = title

    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Convert slide PNG images to PPTX")
    parser.add_argument("images_dir", help="Directory containing PNG slide images")
    parser.add_argument("--prompts", dest="prompts_dir", help="Directory containing prompt .md files (for speaker notes)")
    parser.add_argument("--output", dest="output_path", help="Output PPTX path (default: <parent>/presentation.pptx)")
    parser.add_argument("--title", default="Presentation", help="Presentation title metadata")
    args = parser.parse_args()

    if not os.path.isdir(args.images_dir):
        print(f"ERROR: Images directory not found: {args.images_dir}")
        sys.exit(1)

    result = build_pptx(
        images_dir=args.images_dir,
        prompts_dir=args.prompts_dir,
        output_path=args.output_path,
        title=args.title,
    )
    print(f"PPTX created: {result}")
    print(f"Slides: {len(glob.glob(os.path.join(args.images_dir, '*.png')))}")


if __name__ == "__main__":
    main()
